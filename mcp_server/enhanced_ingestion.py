"""
Enhanced Document Processing Pipeline

Advanced document ingestion with support for complex structures,
multiple formats, and robust error handling.
"""

import os
import re
import json
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass, field
from enum import Enum

import chardet
import fitz  # PyMuPDF
from openpyxl import load_workbook
from pptx import Presentation
from bs4 import BeautifulSoup
import html2text
from loguru import logger
import tiktoken

from .ingestion import DocumentChunk, DocumentProcessor


class ContentType(Enum):
    """Enhanced content type classification."""
    TEXT = "text"
    CODE = "code"
    CONFIG = "config"
    COMMAND = "command"
    TABLE = "table"
    DIAGRAM = "diagram"
    METADATA = "metadata"
    SECTION_HEADER = "section_header"
    LIST = "list"
    QUOTE = "quote"


@dataclass
class EnhancedDocumentChunk(DocumentChunk):
    """Enhanced document chunk with additional metadata."""
    content_type: ContentType = ContentType.TEXT
    section_hierarchy: List[str] = field(default_factory=list)
    extracted_entities: Dict[str, List[str]] = field(default_factory=dict)
    confidence_score: float = 1.0
    parent_chunk_id: Optional[str] = None
    child_chunk_ids: List[str] = field(default_factory=list)


class EnhancedDocumentProcessor(DocumentProcessor):
    """Enhanced document processor with advanced capabilities."""
    
    def __init__(self, max_chunk_size: int = 4000, overlap_size: int = 200):
        super().__init__(max_chunk_size)
        self.overlap_size = overlap_size
        self.encoding = tiktoken.get_encoding("cl100k_base")
        
        # Content extractors registry
        self.extractors = {
            '.pdf': self._extract_pdf_advanced,
            '.docx': self._extract_docx_advanced,
            '.pptx': self._extract_pptx,
            '.xlsx': self._extract_xlsx,
            '.html': self._extract_html,
            '.htm': self._extract_html,
            '.md': self._extract_markdown_advanced,
            '.txt': self._extract_text_advanced,
        }
    
    def process_file_enhanced(self, file_path: str) -> List[EnhancedDocumentChunk]:
        """Process file with enhanced capabilities."""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        logger.info(f"Enhanced processing of file: {file_path}")
        
        # Extract content with metadata
        content, metadata = self._extract_content_with_metadata(file_path)
        
        # Process and chunk with enhanced features
        return self._chunk_content_enhanced(content, str(file_path), metadata)
    
    def _extract_content_with_metadata(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content with comprehensive metadata."""
        file_extension = file_path.suffix.lower()
        metadata = {
            'file_path': str(file_path),
            'file_size': file_path.stat().st_size,
            'file_extension': file_extension,
            'extraction_method': 'unknown'
        }
        
        # Try specific extractor first
        if file_extension in self.extractors:
            try:
                content, extra_metadata = self.extractors[file_extension](file_path)
                metadata.update(extra_metadata)
                metadata['extraction_method'] = f'specific_{file_extension[1:]}'
                return content, metadata
            except Exception as e:
                logger.warning(f"Specific extraction failed for {file_extension}: {e}")
        
        # Fallback to text extraction
        try:
            content = self._extract_text_with_encoding_detection(file_path)
            metadata['extraction_method'] = 'text_fallback'
            return content, metadata
        except Exception as e:
            logger.error(f"All extraction methods failed: {e}")
            raise
    
    def _extract_pdf_advanced(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Advanced PDF extraction with structure preservation."""
        try:
            doc = fitz.open(file_path)
            content_parts = []
            metadata = {
                'page_count': len(doc),
                'has_images': False,
                'has_tables': False,
                'pdf_metadata': doc.metadata
            }
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                
                # Extract text with formatting
                text = page.get_text()
                if text.strip():
                    content_parts.append(f"--- Page {page_num + 1} ---\n{text}")
                
                # Check for images
                if page.get_images():
                    metadata['has_images'] = True
                    content_parts.append(f"[IMAGE DETECTED ON PAGE {page_num + 1}]")
                
                # Extract tables (basic detection)
                tables = page.find_tables()
                if tables:
                    metadata['has_tables'] = True
                    for table in tables:
                        try:
                            table_data = table.extract()
                            table_text = self._format_table_as_text(table_data)
                            content_parts.append(f"[TABLE ON PAGE {page_num + 1}]\n{table_text}")
                        except Exception as e:
                            logger.warning(f"Failed to extract table: {e}")
            
            doc.close()
            return '\n\n'.join(content_parts), metadata
            
        except Exception as e:
            logger.error(f"Advanced PDF extraction failed: {e}")
            # Fallback to basic PDF extraction
            return super()._extract_pdf(file_path), {'extraction_method': 'basic_pdf_fallback'}
    
    def _extract_docx_advanced(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Advanced DOCX extraction with structure preservation."""
        try:
            import docx
            doc = docx.Document(file_path)
            content_parts = []
            metadata = {
                'paragraph_count': 0,
                'table_count': 0,
                'image_count': 0,
                'style_info': {}
            }
            
            # Extract paragraphs with style information
            for para in doc.paragraphs:
                if para.text.strip():
                    style_name = para.style.name if para.style else 'Normal'
                    metadata['style_info'][style_name] = metadata['style_info'].get(style_name, 0) + 1
                    
                    # Add style markers for headers
                    if 'Heading' in style_name:
                        level = self._extract_heading_level(style_name)
                        content_parts.append(f"{'#' * level} {para.text}")
                    else:
                        content_parts.append(para.text)
                    
                    metadata['paragraph_count'] += 1
            
            # Extract tables
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                
                table_text = self._format_table_as_text(table_data)
                content_parts.append(f"[TABLE]\n{table_text}")
                metadata['table_count'] += 1
            
            return '\n\n'.join(content_parts), metadata
            
        except Exception as e:
            logger.error(f"Advanced DOCX extraction failed: {e}")
            # Fallback to text extraction
            return self._extract_text_with_encoding_detection(file_path), {'extraction_method': 'text_fallback'}
    
    def _extract_pptx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from PowerPoint presentations."""
        try:
            prs = Presentation(file_path)
            content_parts = []
            metadata = {
                'slide_count': len(prs.slides),
                'layout_info': {}
            }
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = [f"--- Slide {slide_num} ---"]
                
                # Track layout usage
                layout_name = slide.slide_layout.name
                metadata['layout_info'][layout_name] = metadata['layout_info'].get(layout_name, 0) + 1
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content.append(shape.text)
                
                content_parts.append('\n'.join(slide_content))
            
            return '\n\n'.join(content_parts), metadata
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {e}")
            raise
    
    def _extract_xlsx(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from Excel spreadsheets."""
        try:
            workbook = load_workbook(file_path, read_only=True)
            content_parts = []
            metadata = {
                'sheet_count': len(workbook.sheetnames),
                'sheet_names': workbook.sheetnames
            }
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                sheet_content = [f"--- Sheet: {sheet_name} ---"]
                
                # Extract data from cells
                for row in sheet.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        row_text = '\t'.join(str(cell) if cell is not None else '' for cell in row)
                        sheet_content.append(row_text)
                
                content_parts.append('\n'.join(sheet_content))
            
            workbook.close()
            return '\n\n'.join(content_parts), metadata
            
        except Exception as e:
            logger.error(f"XLSX extraction failed: {e}")
            raise
    
    def _extract_html(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Extract content from HTML files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                html_content = file.read()
            
            soup = BeautifulSoup(html_content, 'html.parser')
            
            # Convert to markdown for better structure preservation
            h = html2text.HTML2Text()
            h.ignore_links = False
            h.ignore_images = False
            markdown_content = h.handle(html_content)
            
            metadata = {
                'title': soup.title.string if soup.title else None,
                'has_scripts': bool(soup.find_all('script')),
                'has_styles': bool(soup.find_all('style')),
                'link_count': len(soup.find_all('a')),
                'image_count': len(soup.find_all('img'))
            }
            
            return markdown_content, metadata
            
        except Exception as e:
            logger.error(f"HTML extraction failed: {e}")
            raise

    def _extract_markdown_advanced(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Advanced markdown extraction with structure analysis."""
        try:
            content = self._extract_text_with_encoding_detection(file_path)

            # Analyze markdown structure
            metadata = {
                'heading_count': len(re.findall(r'^#+\s', content, re.MULTILINE)),
                'code_block_count': len(re.findall(r'```', content)) // 2,
                'link_count': len(re.findall(r'\[.*?\]\(.*?\)', content)),
                'image_count': len(re.findall(r'!\[.*?\]\(.*?\)', content)),
                'table_count': len(re.findall(r'\|.*\|', content, re.MULTILINE))
            }

            return content, metadata

        except Exception as e:
            logger.error(f"Advanced markdown extraction failed: {e}")
            raise

    def _extract_text_advanced(self, file_path: Path) -> Tuple[str, Dict[str, Any]]:
        """Advanced text extraction with content analysis."""
        try:
            content = self._extract_text_with_encoding_detection(file_path)

            # Analyze text structure
            lines = content.split('\n')
            metadata = {
                'line_count': len(lines),
                'word_count': len(content.split()),
                'char_count': len(content),
                'empty_line_count': sum(1 for line in lines if not line.strip()),
                'avg_line_length': sum(len(line) for line in lines) / len(lines) if lines else 0
            }

            return content, metadata

        except Exception as e:
            logger.error(f"Advanced text extraction failed: {e}")
            raise

    def _format_table_as_text(self, table_data: List[List[str]]) -> str:
        """Format table data as readable text."""
        if not table_data:
            return ""

        # Calculate column widths
        col_widths = []
        for col_idx in range(len(table_data[0])):
            max_width = max(len(str(row[col_idx])) if col_idx < len(row) else 0
                          for row in table_data)
            col_widths.append(max_width)

        # Format rows
        formatted_rows = []
        for row in table_data:
            formatted_cells = []
            for col_idx, cell in enumerate(row):
                if col_idx < len(col_widths):
                    formatted_cells.append(str(cell).ljust(col_widths[col_idx]))
                else:
                    formatted_cells.append(str(cell))
            formatted_rows.append(' | '.join(formatted_cells))

        return '\n'.join(formatted_rows)

    def _extract_heading_level(self, style_name: str) -> int:
        """Extract heading level from style name."""
        match = re.search(r'(\d+)', style_name)
        return int(match.group(1)) if match else 1

    def _chunk_content_enhanced(self, content: str, source_file: str,
                              file_metadata: Dict[str, Any]) -> List[EnhancedDocumentChunk]:
        """Enhanced content chunking with semantic awareness."""
        chunks = []

        # Split content into semantic sections
        sections = self._split_into_semantic_sections(content)

        chunk_index = 0
        for section in sections:
            # Analyze section content
            content_type = self._classify_content_type_enhanced(section)
            entities = self._extract_entities(section)
            hierarchy = self._extract_section_hierarchy(section)

            # Handle large sections with overlap
            if self._count_tokens(section) > self.max_chunk_size:
                sub_chunks = self._split_with_overlap(section)
                for i, sub_chunk in enumerate(sub_chunks):
                    chunks.append(EnhancedDocumentChunk(
                        content=sub_chunk,
                        metadata={
                            'token_count': self._count_tokens(sub_chunk),
                            'file_metadata': file_metadata,
                            'section_index': chunk_index,
                            'sub_chunk_index': i,
                            'total_sub_chunks': len(sub_chunks)
                        },
                        chunk_index=chunk_index,
                        source_file=source_file,
                        chunk_type=content_type.value,
                        content_type=content_type,
                        section_hierarchy=hierarchy,
                        extracted_entities=entities,
                        confidence_score=self._calculate_confidence_score(sub_chunk, content_type)
                    ))
                    chunk_index += 1
            else:
                chunks.append(EnhancedDocumentChunk(
                    content=section,
                    metadata={
                        'token_count': self._count_tokens(section),
                        'file_metadata': file_metadata,
                        'section_index': chunk_index
                    },
                    chunk_index=chunk_index,
                    source_file=source_file,
                    chunk_type=content_type.value,
                    content_type=content_type,
                    section_hierarchy=hierarchy,
                    extracted_entities=entities,
                    confidence_score=self._calculate_confidence_score(section, content_type)
                ))
                chunk_index += 1

        logger.info(f"Created {len(chunks)} enhanced chunks from {source_file}")
        return chunks

    def _split_into_semantic_sections(self, content: str) -> List[str]:
        """Split content into semantically meaningful sections."""
        # Enhanced section splitting patterns
        section_patterns = [
            r'\n#{1,6}\s+.*\n',  # Markdown headers
            r'\n={3,}\n',        # Separator lines
            r'\nTask \d+[:\.]',   # Task sections
            r'\nPhase \d+[:\.]',  # Phase sections
            r'\nStep \d+[:\.]',   # Step sections
            r'\n\[.*?\]\n',      # Bracketed sections
            r'\n---+\n',         # Horizontal rules
            r'\n\n(?=[A-Z][^a-z]*:)', # All-caps headers with colons
        ]

        sections = [content]
        for pattern in section_patterns:
            new_sections = []
            for section in sections:
                parts = re.split(pattern, section, flags=re.MULTILINE)
                new_sections.extend([part.strip() for part in parts if part.strip()])
            sections = new_sections

        return sections

    def _classify_content_type_enhanced(self, content: str) -> ContentType:
        """Enhanced content type classification."""
        content_lower = content.lower().strip()

        # Section headers
        if re.match(r'^#+\s', content) or re.match(r'^[A-Z][^a-z]*:?\s*$', content, re.MULTILINE):
            return ContentType.SECTION_HEADER

        # Code blocks
        if ('```' in content or content.startswith('//') or content.startswith('#') or
            re.search(r'^\s*(def|class|function|import|from|const|let|var)\s', content, re.MULTILINE)):
            return ContentType.CODE

        # Configuration files
        if (any(ext in content_lower for ext in ['.json', '.yaml', '.yml', '.toml', '.env']) or
            re.search(r'^\s*[\w-]+\s*[:=]', content, re.MULTILINE)):
            return ContentType.CONFIG

        # Commands
        if any(cmd in content_lower for cmd in ['npm install', 'npx', 'cd ', 'mkdir', 'git ', 'pip install']):
            return ContentType.COMMAND

        # Tables
        if '|' in content and re.search(r'\|.*\|.*\|', content):
            return ContentType.TABLE

        # Lists
        if re.search(r'^\s*[-*+]\s', content, re.MULTILINE) or re.search(r'^\s*\d+\.\s', content, re.MULTILINE):
            return ContentType.LIST

        # Quotes
        if content.startswith('>') or re.search(r'^\s*>', content, re.MULTILINE):
            return ContentType.QUOTE

        # Diagrams/ASCII art
        if re.search(r'[┌┐└┘├┤┬┴┼│─]', content) or re.search(r'[+\-|]{3,}', content):
            return ContentType.DIAGRAM

        return ContentType.TEXT

    def _extract_entities(self, content: str) -> Dict[str, List[str]]:
        """Extract entities from content."""
        entities = {
            'technologies': [],
            'commands': [],
            'file_paths': [],
            'urls': [],
            'dependencies': []
        }

        # Technology patterns
        tech_patterns = [
            r'\b(react|vue|angular|next\.?js|nuxt|svelte|typescript|javascript|python|java|go|rust)\b',
            r'\b(docker|kubernetes|aws|azure|gcp|mongodb|postgresql|mysql|redis)\b',
            r'\b(express|fastapi|django|flask|spring|laravel|rails)\b'
        ]

        for pattern in tech_patterns:
            matches = re.findall(pattern, content, re.IGNORECASE)
            entities['technologies'].extend(matches)

        # Commands
        command_patterns = [
            r'\b(npm|yarn|pip|docker|git|kubectl)\s+\w+',
            r'\bnpx\s+[\w@/-]+',
        ]

        for pattern in command_patterns:
            matches = re.findall(pattern, content)
            entities['commands'].extend(matches)

        # File paths
        file_patterns = [
            r'[./][\w/-]+\.\w+',
            r'\b\w+/[\w/-]+',
        ]

        for pattern in file_patterns:
            matches = re.findall(pattern, content)
            entities['file_paths'].extend(matches)

        # URLs
        url_pattern = r'https?://[^\s]+'
        entities['urls'] = re.findall(url_pattern, content)

        # Dependencies (package names)
        dep_patterns = [
            r'(?:npm install|yarn add|pip install)\s+([\w@/-]+)',
            r'"([\w@/-]+)"\s*:\s*"[^"]*"',  # package.json style
        ]

        for pattern in dep_patterns:
            matches = re.findall(pattern, content)
            entities['dependencies'].extend(matches)

        # Remove duplicates and clean up
        for key in entities:
            entities[key] = list(set(entities[key]))

        return entities

    def _extract_section_hierarchy(self, content: str) -> List[str]:
        """Extract section hierarchy from content."""
        hierarchy = []

        # Look for markdown headers
        header_matches = re.findall(r'^(#+)\s+(.+)$', content, re.MULTILINE)
        for level_markers, title in header_matches:
            level = len(level_markers)
            hierarchy.append(f"H{level}: {title.strip()}")

        # Look for other section indicators
        if re.search(r'^Task \d+', content, re.MULTILINE):
            task_match = re.search(r'^(Task \d+[^:\n]*)', content, re.MULTILINE)
            if task_match:
                hierarchy.append(f"Task: {task_match.group(1)}")

        if re.search(r'^Phase \d+', content, re.MULTILINE):
            phase_match = re.search(r'^(Phase \d+[^:\n]*)', content, re.MULTILINE)
            if phase_match:
                hierarchy.append(f"Phase: {phase_match.group(1)}")

        return hierarchy

    def _calculate_confidence_score(self, content: str, content_type: ContentType) -> float:
        """Calculate confidence score for content classification."""
        base_score = 0.8

        # Adjust based on content length
        if len(content) < 50:
            base_score -= 0.2
        elif len(content) > 1000:
            base_score += 0.1

        # Adjust based on content type indicators
        if content_type == ContentType.CODE and ('def ' in content or 'function ' in content):
            base_score += 0.2
        elif content_type == ContentType.COMMAND and any(cmd in content for cmd in ['npm', 'pip', 'git']):
            base_score += 0.2
        elif content_type == ContentType.CONFIG and ('{' in content or ':' in content):
            base_score += 0.1

        return min(1.0, max(0.1, base_score))

    def _split_with_overlap(self, content: str) -> List[str]:
        """Split large content with overlap to preserve context."""
        chunks = []
        words = content.split()

        if not words:
            return [content]

        # Calculate words per chunk based on token estimate
        words_per_token = 0.75  # Rough estimate
        max_words = int(self.max_chunk_size * words_per_token)
        overlap_words = int(self.overlap_size * words_per_token)

        start = 0
        while start < len(words):
            end = min(start + max_words, len(words))
            chunk_words = words[start:end]
            chunks.append(' '.join(chunk_words))

            if end >= len(words):
                break

            # Move start forward, but with overlap
            start = end - overlap_words

        return chunks
